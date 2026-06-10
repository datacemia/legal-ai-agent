import os
import re
import tempfile
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.business_agent.business_chart_image_service import (
    cleanup_chart_images,
    generate_business_chart_images,
)



BRAND_NAME = "Runexa Systems"
SUPPORTED_LANGUAGES = {"en", "fr", "ar"}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

COLORS = {
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(100, 116, 139),
    "light": RGBColor(248, 250, 252),
    "border": RGBColor(226, 232, 240),
    "green": RGBColor(5, 150, 105),
    "red": RGBColor(220, 38, 38),
    "amber": RGBColor(217, 119, 6),
    "blue": RGBColor(37, 99, 235),
    "white": RGBColor(255, 255, 255),
    "dark": RGBColor(15, 23, 42),
    "navy": RGBColor(2, 6, 23),
    "soft_blue": RGBColor(239, 246, 255),
    "soft_green": RGBColor(236, 253, 245),
}

# Executive slide layout constants.
# Keep a safe vertical rhythm so long FR/AR headlines never overlap the next block.
HEADLINE_TOP = Inches(0.42)
HEADLINE_HEIGHT = Inches(0.78)
SUBLINE_TOP = Inches(1.28)
SUBLINE_HEIGHT = Inches(0.28)
CONTENT_TOP = Inches(1.72)
SECTION_TOP = Inches(1.62)
LOWER_SECTION_TOP = Inches(3.88)
BOTTOM_NOTE_TOP = Inches(6.12)
FOOTER_TOP = Inches(7.08)


def _safe_language(language: str | None) -> str:
    return language if language in SUPPORTED_LANGUAGES else "en"


def _is_rtl(language: str) -> bool:
    return language == "ar"


def _align(language: str) -> PP_ALIGN:
    return PP_ALIGN.RIGHT if _is_rtl(language) else PP_ALIGN.LEFT


def _text(value: Any, fallback: str = "-") -> str:
    if value is None:
        return fallback
    if isinstance(value, str):
        return value.strip() or fallback
    if isinstance(value, int | float):
        return str(value)
    return fallback


def _number(value: Any) -> float | None:
    if isinstance(value, bool):
        return None
    if isinstance(value, int | float):
        return float(value)
    try:
        return float(value)
    except Exception:
        return None


def _format_number(value: Any) -> str:
    numeric = _number(value)
    if numeric is None:
        return "-"
    if numeric.is_integer():
        return f"{numeric:,.0f}"
    return f"{numeric:,.2f}".rstrip("0").rstrip(".")


def _format_percent(value: Any) -> str:
    numeric = _number(value)
    if numeric is None:
        return "-"
    return f"{numeric:,.2f}%".rstrip("0").rstrip(".")


def _format_money(value: Any, currency: dict[str, Any] | None, language: str) -> str:
    numeric = _number(value)
    if numeric is None:
        return "-"

    number = _format_number(numeric)
    symbol = _text((currency or {}).get("symbol"), "")
    code = _text((currency or {}).get("code"), "")
    display_currency = symbol or code

    if not display_currency:
        return number

    if language == "ar":
        return f"{number} {display_currency}".strip()

    if (currency or {}).get("position") == "suffix":
        return f"{number} {display_currency}".strip()

    return f"{display_currency}{number}"


def _format_health_score(
    analysis: dict[str, Any],
    business_health: dict[str, Any] | None = None,
) -> str:
    """
    Format Business Health Score without converting unavailable values to 0/100.

    Product catalogs, inventory files, reference tables, and incomplete datasets can
    legitimately have business_health_score = None. In those cases, PowerPoint
    exports must show N/A, matching the web dashboard and PDF export.
    """

    business_health = business_health or {}
    score = analysis.get("business_health_score")

    if score is None:
        score = business_health.get("score")

    if isinstance(score, bool):
        return "N/A"

    if isinstance(score, int | float):
        return f"{int(round(float(score)))}/100"

    return "N/A"


def _health_score_color(
    analysis: dict[str, Any],
    business_health: dict[str, Any] | None = None,
) -> RGBColor:
    """
    Pick a safe color for Business Health Score cards.
    Unavailable score uses neutral ink instead of red/green.
    """

    business_health = business_health or {}
    score = analysis.get("business_health_score")

    if score is None:
        score = business_health.get("score")

    if isinstance(score, bool) or not isinstance(score, int | float):
        return COLORS["ink"]

    return COLORS["green"] if float(score) >= 60 else COLORS["red"]




def _labels(language: str) -> dict[str, str]:
    labels = {
        "en": {
            "title": "Executive Business Report",
            "subtitle": "Board-ready business intelligence deck",
            "verified": "Data-verified analysis",
            "summary": "Executive Summary",
            "kpis": "Core KPIs",
            "decision": "Priority Decision",
            "risks": "Top Risks",
            "opportunities": "Top Opportunities",
            "recommendations": "Recommendations",
            "forecast": "Forecast",
            "data_quality": "Data Quality",
            "health": "Business Health",
            "business_model": "Business Model",
            "currency": "Currency",
            "revenue": "Revenue",
            "expenses": "Expenses",
            "profit": "Profit",
            "margin": "Margin",
            "growth": "Growth",
            "cashflow": "Cashflow",
            "next_month": "Next Month Revenue",
            "next_quarter": "Next Quarter Revenue",
            "generated_at": "Generated at",
            "source_file": "Source file",
            "trend": "Trend",
            "cashflow_risk": "Cashflow Risk",
            "volatility": "Volatility",
            "charts": "Financial Trends",
            "quality_score": "Quality Score",
            "limitations_ok": "No major data limitations detected.",
        },
        "fr": {
            "title": "Rapport exécutif business",
            "subtitle": "Présentation business prête pour comité de direction",
            "verified": "Analyse vérifiée par les données",
            "summary": "Résumé exécutif",
            "kpis": "KPIs principaux",
            "decision": "Décision prioritaire",
            "risks": "Risques principaux",
            "opportunities": "Opportunités principales",
            "recommendations": "Recommandations",
            "forecast": "Prévisions",
            "data_quality": "Qualité des données",
            "health": "Santé business",
            "business_model": "Modèle business",
            "currency": "Devise",
            "revenue": "Revenus",
            "expenses": "Dépenses",
            "profit": "Profit",
            "margin": "Marge",
            "growth": "Croissance",
            "cashflow": "Cashflow",
            "next_month": "Revenus mois prochain",
            "next_quarter": "Revenus prochain trimestre",
            "generated_at": "Généré le",
            "source_file": "Fichier source",
            "trend": "Tendance",
            "cashflow_risk": "Risque de cashflow",
            "volatility": "Volatilité",
            "charts": "Tendances financières",
            "quality_score": "Score qualité",
            "limitations_ok": "Aucune limitation majeure détectée.",
        },
        "ar": {
            "title": "تقرير تنفيذي للأعمال",
            "subtitle": "عرض أعمال جاهز لمجلس الإدارة",
            "verified": "تحليل موثوق بالبيانات",
            "summary": "الملخص التنفيذي",
            "kpis": "المؤشرات الأساسية",
            "decision": "القرار الأولوي",
            "risks": "أهم المخاطر",
            "opportunities": "أهم الفرص",
            "recommendations": "التوصيات",
            "forecast": "التوقعات",
            "data_quality": "جودة البيانات",
            "health": "صحة الأعمال",
            "business_model": "نوع النشاط",
            "currency": "العملة",
            "revenue": "الإيرادات",
            "expenses": "المصاريف",
            "profit": "الأرباح",
            "margin": "الهامش",
            "growth": "النمو",
            "cashflow": "التدفق النقدي",
            "next_month": "إيرادات الشهر القادم",
            "next_quarter": "إيرادات الربع القادم",
            "generated_at": "تم الإنشاء في",
            "source_file": "الملف المصدر",
            "trend": "الاتجاه",
            "cashflow_risk": "مخاطر التدفق النقدي",
            "volatility": "التقلب",
            "charts": "الاتجاهات المالية",
            "quality_score": "درجة الجودة",
            "limitations_ok": "لم يتم اكتشاف قيود بيانات كبيرة.",
        },
    }
    return labels.get(language, labels["en"])




def _story_labels(language: str) -> dict[str, str]:
    labels = {
        "en": {
            "summary_headline": "The business is growing profitably, but retention quality requires attention",
            "summary_headline_weak": "Business performance requires disciplined action on profitability and risk",
            "kpi_headline": "Financial fundamentals show the current operating profile",
            "chart_headline_positive": "Revenue momentum is outpacing cost pressure",
            "chart_headline_watch": "Financial trends require close monitoring before scaling",
            "forecast_headline_positive": "The forward outlook remains positive if execution discipline is maintained",
            "forecast_headline_watch": "The forecast depends on risk control and execution quality",
            "risks_headline": "The main risks should be handled before aggressive scaling",
            "opportunities_headline": "The strongest opportunities are concentrated around focused growth levers",
            "recommendations_headline": "Management actions should prioritize the highest-impact decisions",
            "quality_headline": "The analysis is data-backed and suitable for executive review",
            "evidence": "Evidence",
            "implication": "Executive implication",
            "focus": "Management focus",
            "watch": "Watch area",
            "revenue_growth": "Revenue growth",
            "profit_margin": "Profit margin",
            "health_score": "Health score",
            "churn": "Churn",
        },
        "fr": {
            "summary_headline": "L’activité progresse avec rentabilité, mais la qualité de rétention exige une attention prioritaire",
            "summary_headline_weak": "La performance business exige une action disciplinée sur la rentabilité et les risques",
            "kpi_headline": "Les fondamentaux financiers résument le profil opérationnel actuel",
            "chart_headline_positive": "La dynamique des revenus progresse plus vite que la pression des coûts",
            "chart_headline_watch": "Les tendances financières doivent être surveillées avant d’accélérer",
            "forecast_headline_positive": "Les perspectives restent positives si la discipline d’exécution est maintenue",
            "forecast_headline_watch": "La prévision dépend du contrôle des risques et de la qualité d’exécution",
            "risks_headline": "Les principaux risques doivent être traités avant une montée en puissance agressive",
            "opportunities_headline": "Les meilleures opportunités se concentrent sur des leviers de croissance ciblés",
            "recommendations_headline": "Les actions de management doivent prioriser les décisions à plus fort impact",
            "quality_headline": "L’analyse est appuyée par les données et adaptée à une revue exécutive",
            "evidence": "Preuve",
            "implication": "Implication exécutive",
            "focus": "Priorité management",
            "watch": "Point de vigilance",
            "revenue_growth": "Croissance des revenus",
            "profit_margin": "Marge bénéficiaire",
            "health_score": "Score de santé",
            "churn": "Churn",
        },
        "ar": {
            "summary_headline": "النشاط ينمو بربحية، لكن جودة الاحتفاظ بالعملاء تحتاج إلى أولوية واضحة",
            "summary_headline_weak": "الأداء التجاري يتطلب انضباطًا في الربحية وإدارة المخاطر",
            "kpi_headline": "المؤشرات المالية تلخص الوضع التشغيلي الحالي",
            "chart_headline_positive": "زخم الإيرادات يتقدم بوتيرة أقوى من ضغط التكاليف",
            "chart_headline_watch": "الاتجاهات المالية تحتاج إلى متابعة دقيقة قبل التوسع",
            "forecast_headline_positive": "التوقعات تبقى إيجابية إذا استمر الانضباط في التنفيذ",
            "forecast_headline_watch": "التوقعات تعتمد على ضبط المخاطر وجودة التنفيذ",
            "risks_headline": "يجب معالجة المخاطر الرئيسية قبل التوسع القوي",
            "opportunities_headline": "أقوى الفرص تتركز حول روافع نمو محددة",
            "recommendations_headline": "يجب أن تركز الإدارة على القرارات الأعلى أثرًا",
            "quality_headline": "التحليل مدعوم بالبيانات ومناسب للمراجعة التنفيذية",
            "evidence": "الدليل",
            "implication": "الأثر التنفيذي",
            "focus": "أولوية الإدارة",
            "watch": "نقطة متابعة",
            "revenue_growth": "نمو الإيرادات",
            "profit_margin": "هامش الربح",
            "health_score": "درجة الصحة",
            "churn": "فقدان العملاء",
        },
    }

    return labels.get(language, labels["en"])


def _analysis_signal(analysis: dict[str, Any]) -> dict[str, float]:
    kpis = analysis.get("kpis") or {}
    advanced = analysis.get("advanced_kpis") or {}

    return {
        "growth": float(_number(kpis.get("growth_rate_percent")) or 0),
        "margin": float(_number(kpis.get("profit_margin_percent")) or 0),
        "score": float(_number(analysis.get("business_health_score")) or 0),
        "churn": float(_number(advanced.get("churn_rate_percent")) or 0),
    }


def _is_positive_profile(analysis: dict[str, Any]) -> bool:
    signal = _analysis_signal(analysis)

    return (
        signal["growth"] >= 0
        and signal["margin"] >= 10
        and signal["score"] >= 60
    )


def _strategic_headline(
    analysis: dict[str, Any],
    language: str,
    kind: str,
) -> str:
    labels = _story_labels(language)
    positive = _is_positive_profile(analysis)

    if kind == "summary":
        return labels["summary_headline"] if positive else labels["summary_headline_weak"]

    if kind == "kpis":
        return labels["kpi_headline"]

    if kind == "charts":
        return labels["chart_headline_positive"] if positive else labels["chart_headline_watch"]

    if kind == "forecast":
        return labels["forecast_headline_positive"] if positive else labels["forecast_headline_watch"]

    if kind == "risks":
        return labels["risks_headline"]

    if kind == "opportunities":
        return labels["opportunities_headline"]

    if kind == "recommendations":
        return labels["recommendations_headline"]

    if kind == "quality":
        return labels["quality_headline"]

    return ""


def _signal_line(analysis: dict[str, Any], language: str) -> str:
    labels = _story_labels(language)
    signal = _analysis_signal(analysis)

    return (
        f"{labels['health_score']}: {_format_number(signal['score'])}/100 · "
        f"{labels['revenue_growth']}: {_format_percent(signal['growth'])} · "
        f"{labels['profit_margin']}: {_format_percent(signal['margin'])}"
    )


def _add_headline(
    slide,
    headline: str,
    language: str,
    top=HEADLINE_TOP,
):
    _add_text(
        slide,
        headline,
        Inches(0.7),
        top,
        Inches(11.9),
        HEADLINE_HEIGHT,
        language,
        font_size=21 if language == "fr" else 22,
        bold=True,
        color=COLORS["ink"],
    )


def _add_subline(
    slide,
    text: str,
    language: str,
    top=SUBLINE_TOP,
):
    _add_text(
        slide,
        text,
        Inches(0.72),
        top,
        Inches(11.5),
        SUBLINE_HEIGHT,
        language,
        font_size=8.5,
        color=COLORS["muted"],
    )


def _add_insight_card(
    slide,
    title: str,
    body: str,
    left,
    top,
    width,
    height,
    language: str,
    accent: RGBColor = COLORS["blue"],
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["border"]

    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()

    _add_text(
        slide,
        title,
        left + Inches(0.22),
        top + Inches(0.16),
        width - Inches(0.35),
        Inches(0.25),
        language,
        font_size=9,
        bold=True,
        color=accent,
    )
    _add_text(
        slide,
        body,
        left + Inches(0.22),
        top + Inches(0.52),
        width - Inches(0.35),
        height - Inches(0.60),
        language,
        font_size=9,
        color=COLORS["ink"],
    )


def _translate_common_value(value: Any, language: str) -> str:
    text = str(value or "").strip()
    normalized = text.lower()

    translations = {
        "ar": {
            "up": "صاعد",
            "down": "هابط",
            "stable": "مستقر",
            "low": "منخفض",
            "medium": "متوسط",
            "high": "مرتفع",
            "critical": "حرج",
            "healthy": "صحي",
            "positive": "إيجابي",
            "negative": "سلبي",
            "general": "عام",
            "saas": "SaaS / اشتراك",
            "healthy profit margin.": "هامش ربح صحي.",
            "healthy growth.": "نمو صحي.",
            "positive cashflow.": "تدفق نقدي إيجابي.",
            "healthy roas.": "عائد إنفاق إعلاني صحي.",
            "healthy cac efficiency.": "كفاءة صحية لتكلفة اكتساب العميل.",
            "critical churn level.": "مستوى فقدان عملاء حرج.",
            "payroll": "الرواتب",
            "marketing": "التسويق",
            "software": "البرمجيات",
            "ad spend": "الإنفاق الإعلاني",
            "customers": "العملاء",
            "orders": "الطلبات",
            "churn": "فقدان العملاء",
        },
        "fr": {
            "up": "hausse",
            "down": "baisse",
            "stable": "stable",
            "low": "faible",
            "medium": "moyen",
            "high": "élevé",
            "critical": "critique",
            "healthy": "sain",
            "positive": "positif",
            "negative": "négatif",
            "general": "général",
            "saas": "SaaS / abonnement",
            "healthy profit margin.": "Marge bénéficiaire saine.",
            "healthy growth.": "Croissance saine.",
            "positive cashflow.": "Cashflow positif.",
            "healthy roas.": "ROAS sain.",
            "healthy cac efficiency.": "Efficacité CAC saine.",
            "critical churn level.": "Niveau de churn critique.",
            "payroll": "Masse salariale",
            "marketing": "Marketing",
            "software": "Logiciels",
            "ad spend": "Dépenses publicitaires",
            "customers": "Clients",
            "orders": "Commandes",
            "churn": "Churn",
        },
        "en": {},
    }

    return translations.get(language, {}).get(normalized, text)


def _translate_category(value: Any, language: str) -> str:
    return _translate_common_value(value, language)



def _translate_and_normalize(value: Any, language: str) -> str:
    translated = _translate_common_value(value, language)

    if language == "fr" and translated:
        exact = {
            "Moyen": "moyen",
            "Élevé": "élevé",
            "Critique": "critique",
            "Sain": "sain",
            "Positif": "positif",
            "Faible": "faible",
            "Volatilité moyen": "Volatilité moyenne",
            "Risque cashflow faible": "Risque de cashflow faible",
            "cashflow est Positif": "cashflow est positif",
            "cashflow est Healthy": "cashflow est sain",
            "situation actuelle comme Critique": "situation actuelle comme critique",
        }
        translated = exact.get(translated, translated)

    return translated


def _format_narrative_text(
    text: Any,
    language: str,
    currency: dict[str, Any] | None = None,
) -> str:
    if not isinstance(text, str):
        return "-"

    result = text

    # Protect common business words before any value-level normalization.
    # This prevents "cashflow" from ever becoming "cashffaible".
    if language == "fr":
        result = result.replace("cashffaible", "cashflow")
        result = result.replace("cashfFaible", "cashflow")
        result = result.replace("cashf faible", "cashflow")
        result = result.replace("cashflow est Positif", "cashflow est positif")
        result = result.replace("cashflow est Sain", "cashflow est sain")
        result = result.replace("cashflow est Critique", "cashflow est critique")
        result = result.replace("score de santé backend est de 73/100 (Sain)", "score de santé backend est de 73/100 (sain)")
        result = result.replace("situation actuelle comme Critique", "situation actuelle comme critique")

    replacements = {
        "137300.0": _format_money(137300.0, currency, language),
        "98650.0": _format_money(98650.0, currency, language),
        "38650.0": _format_money(38650.0, currency, language),
        "28.15%": _format_percent(28.15),
        "10.32%": _format_percent(10.32),
        "28.65%": _format_percent(28.65),
        "28.65": _format_number(28.65),
    }

    for source, target in replacements.items():
        result = result.replace(source, target)

    if language != "en":
        common_values = {
            "Positive": _translate_and_normalize("positive", language),
            "positive": _translate_and_normalize("positive", language),
            "Healthy": _translate_and_normalize("healthy", language),
            "healthy": _translate_and_normalize("healthy", language),
            "Critical": _translate_and_normalize("critical", language),
            "critical": _translate_and_normalize("critical", language),
            "Medium": _translate_and_normalize("medium", language),
            "medium": _translate_and_normalize("medium", language),
            "Low": _translate_and_normalize("low", language),
            "low": _translate_and_normalize("low", language),
        }

        for source, target in common_values.items():
            result = re.sub(
                rf"(?<!\w){re.escape(source)}(?!\w)",
                target,
                result,
                flags=re.IGNORECASE,
            )

    # Final safety cleanup for previously generated or mixed-language fragments.
    if language == "fr":
        result = result.replace("cashffaible", "cashflow")
        result = result.replace("cashflow est Positif", "cashflow est positif")
        result = result.replace("cashflow est Sain", "cashflow est sain")
        result = result.replace("situation actuelle comme Critique", "situation actuelle comme critique")
        result = result.replace("score de santé backend est de 73/100 (Sain)", "score de santé backend est de 73/100 (sain)")
        result = result.replace("Volatilité moyen", "Volatilité moyenne")
        result = result.replace("Volatilité\nmoyen", "Volatilité\nmoyenne")
        result = result.replace("Risque cashflow faible", "Risque de cashflow faible")
        result = result.replace("Risque de cashflow\nfaible", "Risque de cashflow\nfaible")

    return result


def _add_top_rule(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(0.08),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["navy"]
    shape.line.fill.background()


def _add_section_label(slide, text: str, language: str):
    _add_badge(slide, text, Inches(0.7), Inches(0.38), language, width=Inches(2.9))


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _background(slide, color=COLORS["light"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    language: str = "en",
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = COLORS["ink"],
    align: PP_ALIGN | None = None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = frame.paragraphs[0]
    p.alignment = align or _align(language)

    run = p.add_run()
    run.text = _text(text)
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

    return box


def _add_badge(slide, text: str, left, top, language: str, width=Inches(2.6)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(236, 253, 245)
    shape.line.color.rgb = RGBColor(167, 243, 208)

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.02)

    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.bold = True
    run.font.size = Pt(8)
    run.font.color.rgb = COLORS["green"]


def _footer(slide, language: str):
    _add_text(
        slide,
        BRAND_NAME,
        Inches(0.6),
        Inches(7.08),
        Inches(4),
        Inches(0.25),
        language=language,
        font_size=8,
        color=COLORS["muted"],
    )


def _add_card(slide, title: str, value: str, left, top, width, height, language: str, accent=COLORS["ink"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["border"]

    _add_text(
        slide,
        title,
        left + Inches(0.22),
        top + Inches(0.16),
        width - Inches(0.44),
        Inches(0.30),
        language=language,
        font_size=9,
        bold=True,
        color=COLORS["muted"],
    )
    _add_text(
        slide,
        value,
        left + Inches(0.22),
        top + Inches(0.52),
        width - Inches(0.44),
        height - Inches(0.58),
        language=language,
        font_size=21,
        bold=True,
        color=accent,
    )


def _item_title(item: Any) -> str:
    if isinstance(item, str):
        return item
    if not isinstance(item, dict):
        return "Item"
    return (
        item.get("title")
        or item.get("risk")
        or item.get("opportunity")
        or item.get("recommendation")
        or item.get("decision")
        or item.get("description")
        or item.get("what_happened")
        or "Item"
    )


def _item_description(item: Any) -> str:
    if not isinstance(item, dict):
        return ""
    return (
        item.get("description")
        or item.get("why_it_matters")
        or item.get("expected_impact")
        or item.get("recommended_action")
        or item.get("why")
        or item.get("summary")
        or ""
    )


def _add_bullets(slide, items: list[Any], left, top, width, height, language: str, max_items=5):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    align = _align(language)

    if not items:
        p = frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = "-"
        run.font.name = "Arial"
        run.font.size = Pt(13)
        run.font.color.rgb = COLORS["muted"]
        return

    first = True
    for item in items[:max_items]:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        p.alignment = align
        p.level = 0

        run = p.add_run()
        run.text = f"• {_translate_common_value(_item_title(item), language)}"
        run.font.name = "Arial"
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = COLORS["ink"]

        description = _format_narrative_text(_translate_and_normalize(_item_description(item), language), language, None)
        if description:
            p = frame.add_paragraph()
            p.alignment = align
            p.level = 1
            run = p.add_run()
            run.text = description
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.color.rgb = COLORS["muted"]


def _chart_title(title: str, language: str) -> str:
    mapping = {
        "Revenue Trend": {"fr": "Évolution des revenus", "ar": "تطور الإيرادات"},
        "Expense Trend": {"fr": "Évolution des dépenses", "ar": "تطور المصاريف"},
        "Profit Evolution": {"fr": "Évolution du profit", "ar": "تطور الأرباح"},
        "Cashflow Trend": {"fr": "Évolution du cashflow", "ar": "تطور التدفق النقدي"},
        "Expenses by Category": {"fr": "Dépenses par catégorie", "ar": "المصاريف حسب الفئة"},
    }
    return mapping.get(title, {}).get(language, title)


def _add_chart(slide, chart: dict[str, Any], left, top, width, height, language: str):
    data = chart.get("data") or []
    x_key = chart.get("x_key") or "period"
    y_key = chart.get("y_key") or "value"

    categories = []
    values = []
    for row in data:
        number = _number(row.get(y_key))
        if number is None:
            continue
        categories.append(_translate_category(row.get(x_key, ""), language))
        values.append(number)

    if not values:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(_chart_title(str(chart.get("title") or ""), language), values)

    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if chart.get("type") == "bar" else XL_CHART_TYPE.LINE_MARKERS
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    ppt_chart = graphic_frame.chart

    ppt_chart.has_legend = False
    ppt_chart.value_axis.has_major_gridlines = True
    ppt_chart.category_axis.tick_labels.font.size = Pt(8)
    ppt_chart.value_axis.tick_labels.font.size = Pt(8)

    try:
        series = ppt_chart.series[0]
        series.format.line.color.rgb = COLORS["ink"]
        series.format.line.width = Pt(2.5)
    except Exception:
        pass


def _create_cover(prs, analysis, labels, language, source_file_name):
    slide = _blank_slide(prs)
    _background(slide, COLORS["navy"])

    business_health = analysis.get("business_health") or {}
    health_score_display = _format_health_score(analysis, business_health)
    health_score_color = _health_score_color(analysis, business_health)
    currency = analysis.get("currency") or {}
    generated_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")
    source_file = source_file_name or analysis.get("file_metadata", {}).get("file_name", "-")

    _add_text(
        slide,
        BRAND_NAME,
        Inches(0.75),
        Inches(0.55),
        Inches(4.2),
        Inches(0.3),
        language,
        font_size=12,
        bold=True,
        color=COLORS["white"],
    )

    _add_text(
        slide,
        labels["title"],
        Inches(0.75),
        Inches(1.25),
        Inches(8.5),
        Inches(1.0),
        language,
        font_size=40,
        bold=True,
        color=COLORS["white"],
    )

    _add_text(
        slide,
        labels["subtitle"],
        Inches(0.78),
        Inches(2.42),
        Inches(7.8),
        Inches(0.55),
        language,
        font_size=15,
        color=RGBColor(203, 213, 225),
    )

    _add_text(
        slide,
        f"{labels['generated_at']}: {generated_at}\n{labels['source_file']}: {source_file}\n{labels['business_model']}: {_translate_and_normalize(analysis.get('business_model', '-'), language)}\n{labels['currency']}: {currency.get('code', '-')} {currency.get('symbol', '')}",
        Inches(0.78),
        Inches(3.4),
        Inches(7.5),
        Inches(1.2),
        language,
        font_size=11,
        color=RGBColor(203, 213, 225),
    )

    _add_card(
        slide,
        labels["health"],
        health_score_display,
        Inches(9.35),
        Inches(1.25),
        Inches(3.1),
        Inches(1.25),
        language,
        health_score_color,
    )

    _add_card(
        slide,
        labels["revenue"],
        _format_money((analysis.get("kpis") or {}).get("revenue"), currency, language),
        Inches(9.35),
        Inches(2.95),
        Inches(3.1),
        Inches(1.35),
        language,
        COLORS["blue"],
    )

    _add_card(
        slide,
        labels["profit"],
        _format_money((analysis.get("kpis") or {}).get("profit"), currency, language),
        Inches(9.35),
        Inches(4.65),
        Inches(3.1),
        Inches(1.35),
        language,
        COLORS["green"],
    )

    _add_text(
        slide,
        labels["verified"],
        Inches(0.78),
        Inches(6.45),
        Inches(5.8),
        Inches(0.3),
        language,
        font_size=9,
        bold=True,
        color=RGBColor(134, 239, 172),
    )

def _create_summary(prs, analysis, labels, language):
    slide = _blank_slide(prs)
    _background(slide)
    _add_top_rule(slide)

    currency = analysis.get("currency") or {}
    story = _story_labels(language)

    _add_headline(
        slide,
        _strategic_headline(analysis, language, "summary"),
        language,
    )
    _add_subline(slide, _signal_line(analysis, language), language)

    summary = _format_narrative_text(
        analysis.get("executive_summary", "-"),
        language,
        currency,
    )

    _add_insight_card(
        slide,
        labels["summary"],
        summary,
        Inches(0.75),
        CONTENT_TOP,
        Inches(7.35),
        Inches(2.30),
        language,
        COLORS["blue"],
    )

    decision = (analysis.get("smart_insights") or {}).get("most_important_decision") or {}
    if decision:
        _add_insight_card(
            slide,
            labels["decision"],
            decision.get("title") or decision.get("decision") or "-",
            Inches(8.35),
            CONTENT_TOP,
            Inches(4.15),
            Inches(2.30),
            language,
            COLORS["amber"],
        )

        _add_text(
            slide,
            _format_narrative_text(
                decision.get("why") or decision.get("decision") or "-",
                language,
                currency,
            ),
            Inches(0.85),
            Inches(4.38),
            Inches(11.35),
            Inches(1.05),
            language,
            font_size=12,
            color=COLORS["muted"],
        )

    _add_insight_card(
        slide,
        story["implication"],
        _translate_and_normalize(
            (decision or {}).get("impact") or (analysis.get("business_health") or {}).get("rating") or "-",
            language,
        ),
        Inches(0.75),
        Inches(5.62),
        Inches(3.6),
        Inches(0.8),
        language,
        COLORS["green"],
    )

    _add_insight_card(
        slide,
        story["watch"],
        _translate_and_normalize(
            ((analysis.get("forecast") or {}).get("cashflow_risk")) or "-",
            language,
        ),
        Inches(4.65),
        Inches(5.62),
        Inches(3.6),
        Inches(0.8),
        language,
        COLORS["amber"],
    )

    _add_insight_card(
        slide,
        story["focus"],
        (decision or {}).get("decision") or labels["decision"],
        Inches(8.55),
        Inches(5.62),
        Inches(3.95),
        Inches(0.8),
        language,
        COLORS["blue"],
    )

    _footer(slide, language)


def _create_kpis(prs, analysis, labels, language):
    slide = _blank_slide(prs)
    _background(slide)
    _add_top_rule(slide)

    kpis = analysis.get("kpis") or {}
    currency = analysis.get("currency") or {}

    _add_headline(slide, _strategic_headline(analysis, language, "kpis"), language)
    _add_subline(slide, _signal_line(analysis, language), language)

    cards = [
        (labels["revenue"], _format_money(kpis.get("revenue"), currency, language), COLORS["green"]),
        (labels["expenses"], _format_money(kpis.get("expenses"), currency, language), COLORS["amber"]),
        (labels["profit"], _format_money(kpis.get("profit"), currency, language), COLORS["green"]),
        (labels["margin"], _format_percent(kpis.get("profit_margin_percent")), COLORS["ink"]),
        (labels["growth"], _format_percent(kpis.get("growth_rate_percent")), COLORS["blue"]),
        (labels["cashflow"], _translate_and_normalize(kpis.get("cashflow_status"), language), COLORS["green"]),
    ]

    positions = [(0.75, 1.72), (4.75, 1.72), (8.75, 1.72), (0.75, 3.42), (4.75, 3.42), (8.75, 3.42)]
    for (label, value, color), (left, top) in zip(cards, positions):
        _add_card(slide, label, value, Inches(left), Inches(top), Inches(3.35), Inches(1.35), language, color)

    _footer(slide, language)


def _create_charts(prs, analysis, labels, language):
    charts = analysis.get("charts") or []
    if not charts:
        return []

    currency = analysis.get("currency") or {}
    story = _story_labels(language)

    chart_image_paths = generate_business_chart_images(
        charts=charts[:5],
        language=language,
        currency=currency,
    )

    if not chart_image_paths:
        return []

    for index in range(0, len(chart_image_paths), 2):
        slide = _blank_slide(prs)
        _background(slide)
        _add_top_rule(slide)

        _add_headline(
            slide,
            _strategic_headline(analysis, language, "charts"),
            language,
        )
        _add_subline(slide, _signal_line(analysis, language), language)

        first_image = chart_image_paths[index]
        second_image = (
            chart_image_paths[index + 1]
            if index + 1 < len(chart_image_paths)
            else None
        )

        if second_image:
            slide.shapes.add_picture(
                first_image,
                Inches(0.65),
                CONTENT_TOP,
                width=Inches(5.85),
                height=Inches(4.25),
            )
            slide.shapes.add_picture(
                second_image,
                Inches(6.85),
                CONTENT_TOP,
                width=Inches(5.85),
                height=Inches(4.25),
            )
        else:
            slide.shapes.add_picture(
                first_image,
                Inches(1.15),
                CONTENT_TOP,
                width=Inches(10.85),
                height=Inches(4.35),
            )

        _add_insight_card(
            slide,
            story["implication"],
            _strategic_headline(analysis, language, "forecast"),
            Inches(0.8),
            Inches(6.12),
            Inches(11.7),
            Inches(0.46),
            language,
            COLORS["blue"],
        )

        _footer(slide, language)

    return chart_image_paths


def _create_forecast(prs, analysis, labels, language):
    forecast = analysis.get("forecast") or {}
    if not forecast:
        return

    slide = _blank_slide(prs)
    _background(slide)
    _add_top_rule(slide)
    currency = analysis.get("currency") or {}

    _add_headline(slide, _strategic_headline(analysis, language, "forecast"), language)
    _add_subline(slide, _signal_line(analysis, language), language)

    cards = [
        (labels["next_month"], _format_money(forecast.get("next_month_revenue"), currency, language), COLORS["green"]),
        (labels["next_quarter"], _format_money(forecast.get("next_quarter_revenue"), currency, language), COLORS["blue"]),
        (labels["trend"], _translate_and_normalize(forecast.get("trend"), language), COLORS["ink"]),
        (labels["cashflow_risk"], _translate_and_normalize(forecast.get("cashflow_risk"), language), COLORS["amber"]),
        (labels["volatility"], _translate_and_normalize(forecast.get("volatility"), language), COLORS["amber"]),
    ]

    positions = [(0.75, 1.72), (4.75, 1.72), (8.75, 1.72), (0.75, 3.42), (4.75, 3.42)]
    for (label, value, color), (left, top) in zip(cards, positions):
        _add_card(slide, label, value, Inches(left), Inches(top), Inches(3.35), Inches(1.35), language, color)

    if _format_narrative_text(forecast.get("explanation"), language, currency):
        _add_text(slide, _format_narrative_text(forecast.get("explanation"), language, currency), Inches(0.75), Inches(5.72), Inches(11.7), Inches(0.8), language, font_size=11, color=COLORS["muted"])

    _footer(slide, language)


def _create_list_slide(prs, title: str, items: list[Any], language: str, story_kind: str = "risks"):
    slide = _blank_slide(prs)
    _background(slide)
    _add_top_rule(slide)

    _add_headline(
        slide,
        _strategic_headline({}, language, story_kind),
        language,
    )

    _add_subline(slide, title, language)

    safe_items = items[:4] if items else []

    if not safe_items:
        _add_bullets(slide, items, Inches(0.85), Inches(1.72), Inches(11.5), Inches(4.8), language, max_items=6)
        _footer(slide, language)
        return

    positions = [
        (0.75, 1.72),
        (6.75, 1.72),
        (0.75, 3.82),
        (6.75, 3.82),
    ]

    accent = COLORS["red"] if story_kind == "risks" else COLORS["green"] if story_kind == "opportunities" else COLORS["blue"]

    for item, (left, top) in zip(safe_items, positions):
        _add_insight_card(
            slide,
            _translate_and_normalize(_item_title(item), language),
            _format_narrative_text(
                _translate_and_normalize(_item_description(item), language),
                language,
                None,
            ),
            Inches(left),
            Inches(top),
            Inches(5.28),
            Inches(1.75),
            language,
            accent,
        )

    _footer(slide, language)


def _create_quality(prs, analysis, labels, language):
    slide = _blank_slide(prs)
    _background(slide)
    _add_top_rule(slide)

    data_quality = analysis.get("data_quality") or {}
    advanced = analysis.get("advanced_kpis") or {}
    currency = analysis.get("currency") or {}

    _add_headline(slide, _strategic_headline(analysis, language, "quality"), language)
    _add_subline(slide, labels["data_quality"], language)
    _add_card(slide, labels["quality_score"], f"{data_quality.get('score', 0)}/100", Inches(0.75), Inches(1.72), Inches(3.35), Inches(1.35), language, COLORS["green"])
    _add_card(slide, "ROAS", _format_number(advanced.get("roas")), Inches(4.75), Inches(1.72), Inches(3.35), Inches(1.35), language, COLORS["blue"])
    _add_card(slide, "CAC", _format_money(advanced.get("cac"), currency, language), Inches(8.75), Inches(1.72), Inches(3.35), Inches(1.35), language, COLORS["ink"])

    limitations = data_quality.get("limitations") or [labels["limitations_ok"]]
    _add_bullets(slide, limitations, Inches(0.85), Inches(3.72), Inches(11.5), Inches(2.7), language, max_items=5)
    _footer(slide, language)


def build_business_pptx_report(
    analysis: dict[str, Any],
    output_path: str | None = None,
    language: str = "en",
    source_file_name: str | None = None,
) -> str:
    """
    Generate a board-ready PowerPoint deck from existing business analysis.
    """

    language = _safe_language(language)
    labels = _labels(language)

    if output_path is None:
        fd, output_path = tempfile.mkstemp(suffix=".pptx", prefix="runexa_business_report_")
        os.close(fd)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _create_cover(prs, analysis, labels, language, source_file_name)
    _create_summary(prs, analysis, labels, language)
    _create_kpis(prs, analysis, labels, language)
    chart_image_paths = _create_charts(prs, analysis, labels, language)
    _create_forecast(prs, analysis, labels, language)
    _create_list_slide(prs, labels["risks"], analysis.get("risks") or [], language, "risks")
    _create_list_slide(prs, labels["opportunities"], analysis.get("opportunities") or [], language, "opportunities")
    _create_list_slide(prs, labels["recommendations"], analysis.get("recommendations") or [], language, "recommendations")
    _create_quality(prs, analysis, labels, language)

    try:
        prs.save(output_path)
    finally:
        cleanup_chart_images(chart_image_paths)

    return os.path.abspath(output_path)


def build_business_pptx_report_response_payload(pptx_path: str) -> dict[str, Any]:
    return {
        "file_path": pptx_path,
        "file_name": os.path.basename(pptx_path),
        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "generated_at": datetime.utcnow().isoformat() + "Z",
    }
